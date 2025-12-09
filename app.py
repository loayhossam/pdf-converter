
import os
import uuid
import io
from flask import Flask, request, send_file, render_template, after_this_request
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

# إعداد مجلد مؤقت لحفظ الملفات
UPLOAD_FOLDER = '/tmp'
if not os.path.exists(UPLOAD_FOLDER):
    try:
        os.makedirs(UPLOAD_FOLDER)
    except:
        pass

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # الحد الأقصى 16 ميجا

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    # التأكد من وجود ملف
    if 'file' not in request.files:
        return {"error": "لا يوجد ملف"}, 400
    
    file = request.files['file']
    if file.filename == '':
        return {"error": "لم يتم اختيار ملف"}, 400

    if file and file.filename.lower().endswith('.pdf'):
        try:
            # 1. إنشاء اسم عشوائي للملف عشان الملفات متدخلش في بعض
            unique_id = str(uuid.uuid4())
            original_filename = file.filename
            pdf_filename = f"{unique_id}.pdf"
            pptx_filename = f"{unique_id}.pptx"
            
            pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
            pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], pptx_filename)
            
            # حفظ ملف PDF مؤقتاً
            file.save(pdf_path)

            # 2. تحويل PDF إلى صور
            # dpi=150 جودة ممتازة للشاشات وسريعة في التحويل
            images = convert_from_path(pdf_path, dpi=150) 
            
            # 3. إنشاء عرض باوربوينت
            prs = Presentation()
            # ضبط مقاس الشريحة ليكون Widescreen 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for image in images:
                # اختيار تخطيط شريحة فارغة
                blank_slide_layout = prs.slide_layouts[6] 
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # تحويل الصورة لملف في الذاكرة
                image_stream = io.BytesIO()
                image.save(image_stream, format='JPEG', quality=90)
                image_stream.seek(0)
                
                # وضع الصورة لتملأ الشريحة بالكامل
                slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # حفظ ملف الباوربوينت
            prs.save(pptx_path)

            # 4. دالة التنظيف: حذف الملفات من السيرفر بعد ما المستخدم يحملها
            @after_this_request
            def remove_files(response):
                try:
                    if os.path.exists(pdf_path):
                        os.remove(pdf_path)
                    if os.path.exists(pptx_path):
                        os.remove(pptx_path)
                except Exception as e:
                    print(f"Error removing files: {e}")
                return response

            # اسم الملف اللي هيظهر للمستخدم
            download_name = original_filename.replace('.pdf', '_محول.pptx')
            
            # إرسال الملف للتحميل
            return send_file(pptx_path, as_attachment=True, download_name=download_name)

        except Exception as e:
            # في حالة حدوث أي خطأ
            print(f"Error: {e}")
            return {"error": "حدث خطأ أثناء معالجة الملف، تأكد أنه ملف PDF صالح"}, 500
    else:
        return {"error": "الرجاء رفع ملف PDF فقط"}, 400

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=10000)
